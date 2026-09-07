#!/usr/bin/env python3
"""Generate MCSP-232 viva PowerPoint for Google Slides import."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from lxml import etree
from copy import deepcopy

NAVY = RGBColor(0x0A, 0x22, 0x33)
TEAL = RGBColor(0x0E, 0x74, 0x90)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
PAPER = RGBColor(0xF4, 0xF7, 0xFA)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xE0, 0xF7, 0xFB)
LINE = RGBColor(0xDB, 0xE4, 0xEE)
LIGHT_TEAL = RGBColor(0x15, 0x5E, 0x75)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def set_run(run, size=18, bold=False, color=INK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def fill_shape(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill_shape(shape, color)
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size, bold, color, font)
    return box


def add_lines(slide, l, t, w, h, lines, size=18, color=INK, bold_first=False, spacing=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing)
        run = p.add_run()
        run.text = line
        set_run(run, size, bold=(bold_first and i == 0), color=color)
    return box


def chrome(slide, left_label, right_label, page, total=12):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PAPER)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), CYAN)
    add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), NAVY)
    add_text(slide, Inches(0.55), Inches(0.18), Inches(7), Inches(0.32), left_label.upper(), 11, True, TEAL)
    add_text(slide, Inches(9.2), Inches(0.18), Inches(3.5), Inches(0.32), right_label.upper(), 11, True, TEAL, PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.55), Inches(7.12), Inches(6), Inches(0.28), "Bug Tracker System  ·  MCSP-232", 11, False, MUTED)
    add_text(slide, Inches(10.4), Inches(7.12), Inches(2.3), Inches(0.28), f"{page}  /  {total}", 11, False, MUTED, PP_ALIGN.RIGHT)


def heading(slide, text):
    # split "N. Title"
    add_text(slide, Inches(0.55), Inches(0.48), Inches(12.2), Inches(0.55), text, 30, True, NAVY)


def card(slide, l, t, w, h, title, body, accent=True):
    shp = add_rect(slide, l, t, w, h, WHITE, 0.08)
    shp.line.color.rgb = LINE
    shp.line.width = Pt(1)
    if accent:
        bar = add_rect(slide, l, t, Inches(0.08), h, CYAN)
    add_text(slide, l + Inches(0.22), t + Inches(0.12), w - Inches(0.35), Inches(0.32), title, 15, True, TEAL)
    add_text(slide, l + Inches(0.22), t + Inches(0.42), w - Inches(0.35), h - Inches(0.52), body, 14, False, INK)


def bullet_box(slide, l, t, w, h, items, size=18):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "•  " + item
        set_run(run, size, False, INK)


def add_table(slide, l, t, w, h, rows):
    cols = len(rows[0])
    table_shape = slide.shapes.add_table(len(rows), cols, l, t, w, h)
    table = table_shape.table
    widths = [int(w * x) for x in ([0.14, 0.58, 0.28] if cols == 3 else [1 / cols] * cols)]
    if cols == 3:
        table.columns[0].width = int(w * 0.16)
        table.columns[1].width = int(w * 0.56)
        table.columns[2].width = int(w * 0.28)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    set_run(run, 14, bold=(r == 0), color=WHITE if r == 0 else INK)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = NAVY if r == 0 else (RGBColor(0xF8, 0xFA, 0xFC) if r % 2 else WHITE)
    return table_shape


# ---------- SLIDE 1 TITLE ----------
s = prs.slides.add_slide(blank)
chrome(s, "IGNOU  ·  MCA_NEW  ·  MCSP-232", "Project Viva", 1)
pill = add_rect(s, Inches(0.55), Inches(1.55), Inches(2.5), Inches(0.38), SOFT, 0.5)
add_text(s, Inches(0.55), Inches(1.58), Inches(2.5), Inches(0.34), "FINAL SEMESTER PROJECT", 11, True, LIGHT_TEAL, PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(2.1), Inches(12), Inches(0.9), "Bug Tracker System", 44, True, NAVY)
add_text(s, Inches(0.55), Inches(2.95), Inches(12), Inches(0.5), "A web-based platform for complete software bug lifecycle management", 20, False, TEAL)
add_lines(s, Inches(0.55), Inches(3.7), Inches(12), Inches(1.8), [
    "Submitted by:   Rohit Singh          Enrolment No.:  2401012750",
    "Project Guide:  Deepak Singh Kathait",
    "School of Computer and Information Sciences, IGNOU",
    "Indira Gandhi National Open University, New Delhi",
], 18, INK)

# ---------- SLIDE 2 PROBLEM ----------
s = prs.slides.add_slide(blank)
chrome(s, "Need for the system", "MCSP-232", 2)
heading(s, "1.  Problem Statement")
cards = [
    ("Current practice", "Bugs are often tracked in email, chat, or spreadsheets with no single source of truth."),
    ("No enforced workflow", "Testing and validation can be skipped, so defects may be marked resolved incorrectly."),
    ("Poor coordination", "Reporter, tester and developer work with incomplete or outdated information."),
    ("Weak accountability", "No complete audit trail of who changed a bug, when, and why."),
]
positions = [(0.55, 1.2), (6.9, 1.2), (0.55, 3.15), (6.9, 3.15)]
for (title, body), (x, y) in zip(cards, positions):
    card(s, Inches(x), Inches(y), Inches(5.9), Inches(1.75), title, body)
note = add_rect(s, Inches(0.55), Inches(5.15), Inches(12.2), Inches(1.35), SOFT, 0.08)
add_text(s, Inches(0.75), Inches(5.3), Inches(11.8), Inches(1.05),
         "Tools like Jira are powerful but heavy and costly for small / academic teams. This project offers a structured, lightweight alternative with a mandatory seven-phase quality workflow.",
         16, False, LIGHT_TEAL)

# ---------- SLIDE 3 OBJECTIVES ----------
s = prs.slides.add_slide(blank)
chrome(s, "What we set out to build", "MCSP-232", 3)
heading(s, "2.  Objectives")
bullet_box(s, Inches(0.7), Inches(1.25), Inches(12), Inches(5.4), [
    "Automate the bug lifecycle through a seven-phase workflow",
    "Implement role-based access: Bug Raiser, Tester, Developer, Administrator",
    "Enforce mandatory testing before a bug can be closed",
    "Organize bugs by Tool / application with automatic tester–developer assignment",
    "Maintain a full audit log and notify users of status changes",
    "Provide dashboards for counts, priority and phase-wise progress",
], 20)

# ---------- SLIDE 4 PHASES ----------
s = prs.slides.add_slide(blank)
chrome(s, "Core workflow", "MCSP-232", 4)
heading(s, "3.  Seven-Phase Workflow")
phases = [
    ("I", "Report", "Raiser"),
    ("II", "Confirm", "Tester"),
    ("III", "Analyse", "Developer"),
    ("IV", "Fix", "Developer"),
    ("V", "Test", "Tester"),
    ("VI–VII", "Deploy / Close", "Dev / Admin"),
]
for i, (num, name, who) in enumerate(phases):
    x = Inches(0.45 + i * 2.12)
    add_rect(s, x, Inches(1.2), Inches(1.95), Inches(1.05), NAVY if i == 5 else TEAL, 0.1)
    add_text(s, x, Inches(1.28), Inches(1.95), Inches(0.38), f"{num}  {name}", 13, True, WHITE, PP_ALIGN.CENTER)
    add_text(s, x, Inches(1.68), Inches(1.95), Inches(0.4), who, 12, False, RGBColor(0xE0, 0xF2, 0xFE), PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, x + Inches(1.82), Inches(1.42), Inches(0.35), Inches(0.4), "→", 18, True, CYAN, PP_ALIGN.CENTER)

add_table(s, Inches(0.55), Inches(2.5), Inches(12.2), Inches(2.6), [
    ["Phase", "Activity", "Owner"],
    ["I", "Bug Report — details, expected vs actual, attachments", "Bug Raiser"],
    ["II", "Confirmation — reproduce and accept / reject", "Tester"],
    ["III–IV", "Root-cause analysis and code fix", "Developer"],
    ["V–VII", "Final test, deployment notes, closure & lessons", "Tester / Dev / Admin"],
])
note = add_rect(s, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.15), SOFT, 0.08)
add_text(s, Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.8),
         "Quality gate: if Tester or Analyst rejects the bug, it is marked Closed immediately. This enforces QA and is the main unique point of the project.",
         16, False, LIGHT_TEAL)

# ---------- SLIDE 5 STACK ----------
s = prs.slides.add_slide(blank)
chrome(s, "How it is built", "MCSP-232", 5)
heading(s, "4.  Technology Stack & Architecture")
card(s, Inches(0.55), Inches(1.2), Inches(3.95), Inches(1.7), "Frontend", "React.js, React Router, Axios, Tailwind CSS, Recharts")
card(s, Inches(4.7), Inches(1.2), Inches(3.95), Inches(1.7), "Backend", "Node.js, Express.js, REST APIs, TypeScript")
card(s, Inches(8.85), Inches(1.2), Inches(3.95), Inches(1.7), "Data & security", "MongoDB, Mongoose, JWT, bcrypt, Cloudinary, Nodemailer")

layers = [
    (NAVY, "Presentation Layer", "React.js  ·  UI, routing, client validation  ·  Browser"),
    (TEAL, "Application Layer", "Express.js REST APIs  ·  RBAC middleware  ·  business rules  ·  notifications"),
    (LIGHT_TEAL, "Data Layer", "MongoDB + Mongoose  ·  Users, Tools, Bugs, Config, Logs  ·  file storage"),
]
for i, (col, title, body) in enumerate(layers):
    y = Inches(3.2 + i * 1.1)
    add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.95), col, 0.08)
    add_text(s, Inches(0.75), y + Inches(0.1), Inches(11.8), Inches(0.32), title, 14, True, RGBColor(0xA5, 0xF3, 0xFC), PP_ALIGN.CENTER)
    add_text(s, Inches(0.75), y + Inches(0.42), Inches(11.8), Inches(0.4), body, 16, False, WHITE, PP_ALIGN.CENTER)

# ---------- SLIDE 6 DFD ----------
s = prs.slides.add_slide(blank)
chrome(s, "System analysis", "MCSP-232", 6)
heading(s, "5.  DFD Level-0 (Context Diagram)")
# center process
add_rect(s, Inches(5.15), Inches(3.15), Inches(3.05), Inches(1.35), NAVY, 0.5)
add_text(s, Inches(5.15), Inches(3.4), Inches(3.05), Inches(0.4), "Bug Tracker", 16, True, CYAN, PP_ALIGN.CENTER)
add_text(s, Inches(5.15), Inches(3.8), Inches(3.05), Inches(0.4), "System", 16, False, WHITE, PP_ALIGN.CENTER)

ents = [
    (0.7, 1.25, "Bug Raiser", "Bug report / status"),
    (10.1, 1.25, "Administrator", "Users, settings, reports"),
    (0.7, 5.15, "Tester", "Test results"),
    (10.1, 5.15, "Developer", "Fix / resolution"),
    (5.25, 1.25, "Database", "Store / retrieve"),
]
for x, y, name, flow in ents:
    add_rect(s, Inches(x), Inches(y), Inches(2.4), Inches(0.85), SOFT, 0.08)
    add_text(s, Inches(x), Inches(y + 0.08), Inches(2.4), Inches(0.35), name, 15, True, NAVY, PP_ALIGN.CENTER)
    add_text(s, Inches(x), Inches(y + 0.42), Inches(2.4), Inches(0.35), flow, 12, False, TEAL, PP_ALIGN.CENTER)

add_text(s, Inches(0.55), Inches(2.2), Inches(12.2), Inches(0.7),
         "External entities exchange data with one central process. Database stores and retrieves bug, user and tool records.",
         15, False, MUTED, PP_ALIGN.CENTER)

# ---------- SLIDE 7 ER ----------
s = prs.slides.add_slide(blank)
chrome(s, "Data model", "MCSP-232", 7)
heading(s, "6.  ER Diagram & Collections")
entities = [("USER", 2.15, 1.2), ("TOOL", 0.55, 2.7), ("BUG", 2.15, 2.7), ("LOG", 3.75, 2.7), ("NOTIFICATION", 2.0, 4.2)]
for name, x, y in entities:
    add_rect(s, Inches(x), Inches(y), Inches(1.7), Inches(0.7), NAVY, 0.1)
    add_text(s, Inches(x), Inches(y + 0.16), Inches(1.7), Inches(0.4), name, 13, True, WHITE, PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(1.95), Inches(4.9), Inches(0.55), "User  1:N  Bug     Bug  N:1  Tool     Bug  1:N  Log / Notification", 13, False, TEAL, PP_ALIGN.CENTER)

card(s, Inches(6.5), Inches(1.2), Inches(6.25), Inches(1.0), "Users", "username, email, passwordHash, role, isActive")
card(s, Inches(6.5), Inches(2.3), Inches(6.25), Inches(1.0), "Tools", "toolName, testerId, devId, stack, SOP")
card(s, Inches(6.5), Inches(3.4), Inches(6.25), Inches(1.0), "Bugs", "bugId, currentPhase, 7 embedded phase objects")
card(s, Inches(6.5), Inches(4.5), Inches(6.25), Inches(1.0), "Config + Logs", "stages / status / priority  ·  who did what, when")
note = add_rect(s, Inches(0.55), Inches(5.65), Inches(12.2), Inches(0.9), SOFT, 0.08)
add_text(s, Inches(0.75), Inches(5.8), Inches(11.8), Inches(0.65),
         "MongoDB (NoSQL): phase details are embedded in the bug document so one bug stays together through all 7 stages.",
         15, False, LIGHT_TEAL)

# ---------- SLIDE 8 MODULES ----------
s = prs.slides.add_slide(blank)
chrome(s, "Implementation", "MCSP-232", 8)
heading(s, "7.  Major Modules")
mods = [
    ("1. Authentication", "Register, login, JWT, bcrypt, OTP password reset"),
    ("2. User management", "Roles, activate / deactivate, profile"),
    ("3. Tool management", "Apps, default tester & developer, stack"),
    ("4. Bug workflow", "Phases I–VII with validation at each gate"),
    ("5. Notifications", "In-app alerts for assignment and status"),
    ("6. Dashboard", "Counts, charts, priority and phase stats"),
    ("7. Admin panel", "Users, tools, configuration, logs"),
    ("8. Search & filter", "Multi-criteria list, sort, pagination"),
    ("9. Activity log", "Full audit trail of bug actions"),
]
for i, (title, body) in enumerate(mods):
    r, c = divmod(i, 3)
    card(s, Inches(0.55 + c * 4.2), Inches(1.2 + r * 1.75), Inches(4.0), Inches(1.55), title, body)

# ---------- SLIDE 9 SECURITY ----------
s = prs.slides.add_slide(blank)
chrome(s, "Quality attributes", "MCSP-232", 9)
heading(s, "8.  Security Mechanism")
left = add_rect(s, Inches(0.55), Inches(1.25), Inches(6.0), Inches(5.15), WHITE, 0.08)
left.line.color.rgb = LINE
add_rect(s, Inches(0.55), Inches(1.25), Inches(0.1), Inches(5.15), CYAN)
add_text(s, Inches(0.9), Inches(1.4), Inches(5.4), Inches(0.4), "Authentication", 18, True, TEAL)
bullet_box(s, Inches(0.9), Inches(1.95), Inches(5.4), Inches(4.2), [
    "bcrypt hashing, 10 salt rounds",
    "No plain-text passwords",
    "JWT on protected routes",
    "OTP-based password reset",
    "Login rate limiting",
], 17)
right = add_rect(s, Inches(6.8), Inches(1.25), Inches(6.0), Inches(5.15), WHITE, 0.08)
right.line.color.rgb = LINE
add_rect(s, Inches(6.8), Inches(1.25), Inches(0.1), Inches(5.15), CYAN)
add_text(s, Inches(7.15), Inches(1.4), Inches(5.4), Inches(0.4), "Authorization & data", 18, True, TEAL)
bullet_box(s, Inches(7.15), Inches(1.95), Inches(5.4), Inches(4.2), [
    "RBAC: raiser / tester / developer / admin",
    "Phase-specific rights (assigned tester confirms)",
    "Mongoose validation (NoSQL injection control)",
    "File type & size limits (~10 MB)",
    "Secrets stored in environment variables",
], 17)

# ---------- SLIDE 10 TESTING ----------
s = prs.slides.add_slide(blank)
chrome(s, "Verification", "MCSP-232", 10)
heading(s, "9.  Testing Strategy")
add_table(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(4.0), [
    ["Type", "Scope", "Approach"],
    ["Unit testing", "API functions, components, DB operations", "Jest / isolated checks"],
    ["Integration", "Routes + MongoDB + modules", "API testing"],
    ["System / E2E", "All 7 phases, all four roles", "Manual walkthrough"],
    ["Security", "Auth, RBAC, input validation", "Manual + middleware"],
    ["Performance", "Page and API response", "Target: page < 3s, API < 2s"],
    ["Regression", "Old flow after a new fix", "Smoke test before “done”"],
])
note = add_rect(s, Inches(0.55), Inches(5.45), Inches(12.2), Inches(1.05), SOFT, 0.08)
add_text(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(0.75),
         "Primary viva demo: manual end-to-end test — login → report bug → confirm → analyse → fix → final test → close.",
         16, False, LIGHT_TEAL)

# ---------- SLIDE 11 GANTT ----------
s = prs.slides.add_slide(blank)
chrome(s, "Project planning", "MCSP-232", 11)
heading(s, "10.  Gantt / PERT Highlights")
left = add_rect(s, Inches(0.55), Inches(1.25), Inches(6.0), Inches(4.9), WHITE, 0.08)
left.line.color.rgb = LINE
add_rect(s, Inches(0.55), Inches(1.25), Inches(0.1), Inches(4.9), CYAN)
add_text(s, Inches(0.9), Inches(1.45), Inches(5.4), Inches(0.4), "Milestones (20 weeks)", 16, True, TEAL)
bullet_box(s, Inches(0.9), Inches(2.05), Inches(5.4), Inches(3.8), [
    "Week 3 — Auth & user management",
    "Week 9 — Bug report & confirmation",
    "Week 14 — Full 7-phase workflow",
    "Week 17 — Admin & analytics",
    "Week 20 — Deployment & report",
], 17)
add_rect(s, Inches(6.8), Inches(1.25), Inches(6.0), Inches(4.9), WHITE, 0.08).line.color.rgb = LINE
add_rect(s, Inches(6.8), Inches(1.25), Inches(0.1), Inches(4.9), CYAN)
add_text(s, Inches(7.15), Inches(1.45), Inches(5.4), Inches(0.4), "Critical path", 16, True, TEAL)
add_text(s, Inches(7.15), Inches(2.15), Inches(5.4), Inches(0.7), "1 → 2 → 3 → 4 → 5 → 6 → 7 → 8 → 9", 18, True, NAVY)
bullet_box(s, Inches(7.15), Inches(3.0), Inches(5.4), Inches(2.8), [
    "Longest dependent chain of tasks",
    "Total planned duration: 20 weeks",
    "Full Gantt and PERT charts are in the project report (PDF pages ~70–71)",
], 17)

# ---------- SLIDE 12 END ----------
s = prs.slides.add_slide(blank)
chrome(s, "Wrap-up", "MCSP-232", 12)
heading(s, "11.  Conclusion & Future Scope")
card(s, Inches(0.55), Inches(1.2), Inches(3.95), Inches(2.2), "Achieved", "7-phase tracking, RBAC, logs, MERN implementation, dashboards")
card(s, Inches(4.7), Inches(1.2), Inches(3.95), Inches(2.2), "Limitations", "Web only, no Git/Jira sync, no MFA, English UI only")
card(s, Inches(8.85), Inches(1.2), Inches(3.95), Inches(2.2), "Future", "Mobile app, GitHub integration, AI duplicate-bug detection")
add_text(s, Inches(0.55), Inches(3.7), Inches(12.2), Inches(0.8), "Thank you", 40, True, TEAL, PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(4.5), Inches(12.2), Inches(0.45), "I can demonstrate the live system now.  Questions are welcome.", 18, False, MUTED, PP_ALIGN.CENTER)
add_text(s, Inches(0.55), Inches(5.3), Inches(12.2), Inches(0.8), "Rohit Singh  ·  Enrolment 2401012750  ·  Guide: Deepak Singh Kathait", 16, False, NAVY, PP_ALIGN.CENTER)

out = "/Users/harekrishn/Desktop/Param2/College Project/Bug_Bunters/MCSP232_Viva_Presentation.pptx"
prs.save(out)
print("Saved", out)
